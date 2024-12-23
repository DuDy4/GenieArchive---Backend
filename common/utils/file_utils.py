import os
import io
import docx
import fitz  
import pytesseract
from PIL import Image
from pptx import Presentation

def get_file_name_from_url(url: str) -> str:
    return os.path.basename(url)

def get_file_extension(file_name: str) -> str:
    return os.path.splitext(file_name)[1]
# PDF extraction (using PyMuPDF)
def extract_text_from_pdf2(file_content):
    pdf_document = fitz.open(stream=file_content, filetype="pdf")
    text = ""
    for page in pdf_document:
        text += page.get_text("text")
    return text

def extract_text_from_pdf(file_content):
    try:
        pdf_document = fitz.open(stream=file_content, filetype="pdf")
        text = ""
        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            # Attempt to extract text
            page_text = page.get_text("text")
            if not page_text.strip():  # If no text found, try OCR on the page's image
                pix = page.get_pixmap()
                img = Image.open(io.BytesIO(pix.tobytes()))
                page_text = pytesseract.image_to_string(img)
            text += page_text
        return text
    except Exception as e:
        print(f"An error occurred while extracting text: {e}")
        return ""


def extract_text_from_docx(file_content):
    with io.BytesIO(file_content) as file_like_object:
        doc = docx.Document(file_like_object)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

    if not text:
        file_like_object = io.BytesIO(file_content)
        document = docx.Document(file_like_object)

        # Extract text from paragraphs
        full_text = []
        for paragraph in document.paragraphs:
            full_text.append(paragraph.text)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text)

        return "\n".join(full_text)
    return text


def extract_text_from_pptx(file_content):
    presentation = Presentation(file_content)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
